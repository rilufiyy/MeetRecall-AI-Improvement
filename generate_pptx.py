"""
MeetRecall AI Improvement - Talent Showcase
Non-technical version for HR / hiring partners.
Focus: what the product does + value, not how it's built.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ────────────────────────────────────────────────────────────────────
LIGHT_BLUE = RGBColor(202, 228, 252)
MED_BLUE   = RGBColor( 57,  93, 145)
DARK_NAVY  = RGBColor( 10,  22,  46)
DARK_BG    = RGBColor( 22,  20,  68)
INDIGO_5   = RGBColor( 99, 102, 241)
INDIGO_3   = RGBColor(165, 180, 252)
INDIGO_1   = RGBColor(224, 231, 255)
VIOLET_5   = RGBColor(139,  92, 246)
GREEN_5    = RGBColor( 34, 197,  94)
GREEN_1    = RGBColor(220, 252, 231)
RED_4      = RGBColor(248, 113, 113)
RED_1      = RGBColor(254, 226, 226)
WHITE      = RGBColor(255, 255, 255)
GRAY_7     = RGBColor( 55,  65,  81)
GRAY_5     = RGBColor(107, 114, 128)
GRAY_2     = RGBColor(229, 231, 235)
DARK_CARD  = RGBColor( 28,  25,  85)
AMBER_4    = RGBColor(251, 191,  36)

DOCS      = r"D:\MeetRecall AI Improvement\documentations"
FLOWCHART = r"D:\MeetRecall AI Improvement\documentations\flowchart_new.png"
OUT       = r"D:\MeetRecall AI Improvement\MeetRecall_AI_Improvement_Presentation_v4.pptx"

# ── Primitives ─────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


def slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(s, c):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c


def box(s, l, t, w, h, fill, border=None, bw=1.0):
    sp = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if border:
        sp.line.color.rgb = border; sp.line.width = Pt(bw)
    else:
        sp.line.fill.background()
    return sp


def txt(s, text, l, t, w, h, sz=13, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color


def img(s, fname, l, t, w, h=None):
    path = os.path.join(DOCS, fname)
    if not os.path.exists(path): return
    if h: s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else: s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))


def hdr(s, title, sub=""):
    box(s, 0, 0, 13.33, 1.05, DARK_BG)
    box(s, 0, 1.05, 13.33, 0.05, INDIGO_5)
    txt(s, title, 0.4, 0.13, 9, 0.75, sz=26, bold=True, color=WHITE)
    if sub:
        txt(s, sub, 0.4, 0.72, 12.5, 0.33, sz=10, color=INDIGO_3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 - TITLE
# ══════════════════════════════════════════════════════════════════════════════
s1 = slide()
bg(s1, LIGHT_BLUE)

box(s1, 0.00, 0.76, 4.24, 6.74, MED_BLUE)
box(s1, 2.12, 0.00, 4.24, 4.44, DARK_NAVY)
box(s1, 2.97, 0.00, 4.25, 5.75, MED_BLUE)

txt(s1, "Presented by : Sri Lutfiya Dwiyeni",
    7.4, 0.6, 5.6, 0.52, sz=12, color=DARK_NAVY, align=PP_ALIGN.RIGHT)

txt(s1, "MeetRecall AI\nImprovement",
    7.3, 3.0, 5.8, 2.4, sz=42, bold=True, color=DARK_NAVY)

txt(s1, "Meeting Intelligence Platform",
    7.3, 5.55, 5.8, 0.6, sz=15, color=MED_BLUE)

tags = ["Transcription", "RAG Chatbot", "Topic Clustering", "Sentiment Analysis"]
for i, tag in enumerate(tags):
    c, r = i % 2, i // 2
    box(s1, 7.3 + c * 2.95, 6.28 + r * 0.55, 2.6, 0.42, MED_BLUE)
    txt(s1, tag, 7.35 + c * 2.95, 6.32 + r * 0.55, 2.5, 0.34,
        sz=10, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 - ABOUT ME
# ══════════════════════════════════════════════════════════════════════════════
s2 = slide()
bg(s2, LIGHT_BLUE)

box(s2, 0.00, 0.00, 1.05, 7.50, MED_BLUE)
box(s2, 0.00, 3.95, 1.38, 3.55, DARK_NAVY)
box(s2, 12.28, 0.00, 1.05, 3.39, DARK_NAVY)
box(s2, 12.65, 1.46, 0.68, 6.04, MED_BLUE)

box(s2, 1.38, 3.28, 2.78, 1.57, MED_BLUE)
box(s2, 1.59, 3.56, 1.30, 1.46, DARK_NAVY)
box(s2, 1.59, 4.25, 2.56, 1.57, MED_BLUE)

box(s2, 0.50, 4.56, 3.99, 0.82, DARK_NAVY)
txt(s2, "SRI LUTFIYA DWIYENI",
    0.55, 4.67, 3.88, 0.58, sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s2, "ABOUT ME", 4.80, 0.55, 8.10, 0.72, sz=26, bold=True, color=DARK_NAVY)
box(s2, 4.80, 1.35, 8.10, 0.04, MED_BLUE)

txt(s2, "EDUCATION",  4.80, 1.52, 3.8, 0.4, sz=11, bold=True, color=DARK_NAVY)
txt(s2, "Diponegoro University - Bachelor of Mathematics (2021 - 2025)",
    4.80, 1.92, 7.8, 0.4, sz=11, color=GRAY_7)

txt(s2, "EXPERIENCE", 4.80, 2.48, 3.8, 0.4, sz=11, bold=True, color=DARK_NAVY)
txt(s2, "Independent Study at Bangkit Academy 2024  (Sept 2024 - Jan 2025)",
    4.80, 2.88, 7.8, 0.38, sz=11, color=GRAY_7)
txt(s2, "Independent Study at Startup Campus 2024  (Feb 2024 - Jun 2024)",
    4.80, 3.26, 7.8, 0.38, sz=11, color=GRAY_7)

txt(s2, "OVERVIEW PROJECT", 4.80, 3.76, 8.10, 0.4, sz=11, bold=True, color=DARK_NAVY)

projects = [
    ("Clinical Research RAG Mini Chatbot",
     "Mini chatbot RAG focused on Nightmare Disorder, IRT, and Major Depressive Episode."),
    ("MLOps Model Monitoring",
     "MLOps pipeline for house price prediction using MLflow, FastAPI, and Docker."),
    ("SmartSplit Bill Assistant",
     "AI tool to extract receipts and split bills automatically."),
    ("Intelligent Customer Assistant",
     "RAG-based assistant that answers customer queries accurately and contextually."),
]
for i, (title, desc) in enumerate(projects):
    y = 4.28 + i * 0.75
    txt(s2, title + ":", 4.80, y,        8.1, 0.34, sz=10, bold=True, color=DARK_NAVY)
    txt(s2, desc,         4.80, y + 0.33, 8.1, 0.36, sz=10, color=GRAY_5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 - EXECUTIVE SUMMARY
# Clean, readable format matching reference PDF style
# ══════════════════════════════════════════════════════════════════════════════
s_exec = slide()
bg(s_exec, WHITE)
hdr(s_exec, "Executive Summary")

# Project name + one-liner
txt(s_exec, "MeetRecall AI Improvement",
    0.55, 1.25, 12.2, 0.65, sz=26, bold=True, color=DARK_BG)

txt(s_exec,
    "is a Meeting Intelligence Platform that transforms any meeting recording "
    "into a full set of AI-powered insights — automatically, privately, and "
    "without any manual effort.",
    0.55, 1.9, 12.2, 0.85, sz=15, color=GRAY_7)

# Horizontal rule
box(s_exec, 0.55, 2.85, 12.2, 0.04, INDIGO_5)

# Bullet points — plain language, tailored to the improved version
bullets = [
    (INDIGO_5,
     "Automatically transcribes any meeting recording",
     "Converts speech to text and identifies each speaker — no manual note-taking required."),
    (VIOLET_5,
     "Answers questions about the meeting",
     "Ask anything in plain language and get a direct answer drawn from the actual transcript."),
    (RGBColor(14, 165, 233),
     "Discovers the main topics discussed",
     "The AI reads through the conversation and organises it into clear, labelled themes."),
    (GREEN_5,
     "Reveals how each participant came across",
     "See the emotional tone of every speaker — positive, neutral, or negative — throughout the meeting."),
    (MED_BLUE,
     "Runs privately on your own machine",
     "No data is sent to the cloud. Everything stays on your computer, ready to use with one command."),
]

for i, (color, bold_part, plain_part) in enumerate(bullets):
    y = 3.05 + i * 0.83

    # Colour bar on the left
    box(s_exec, 0.55, y + 0.06, 0.06, 0.58, color)

    # Bold label
    tb = s_exec.shapes.add_textbox(Inches(0.75), Inches(y), Inches(11.8), Inches(0.38))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r_bold = p.add_run()
    r_bold.text = bold_part + "  "
    r_bold.font.size = Pt(13); r_bold.font.bold = True; r_bold.font.color.rgb = DARK_BG

    # Plain continuation on second line
    txt(s_exec, plain_part, 0.75, y + 0.38, 11.8, 0.4, sz=11, color=GRAY_7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 - PROBLEM STATEMENT
# Kept simple and relatable — no jargon
# ══════════════════════════════════════════════════════════════════════════════
s3 = slide()
bg(s3, WHITE)
hdr(s3, "The Challenge",
    "What professionals face every day with meeting recordings")

problems = [
    ("Takes Too Much Time",
     "Writing down everything from a 1-hour meeting can take another 4 to 5 hours of manual work."),
    ("Hard to Find Information",
     "After a meeting, finding out what was said or who decided what requires replaying the entire recording."),
    ("No Clear Picture of Topics",
     "When a meeting covers many subjects, it is nearly impossible to separate and organize each theme manually."),
    ("Missing the Emotional Context",
     "Understanding how each participant felt during the discussion is invisible without the right tool."),
]

for i, (h, b) in enumerate(problems):
    col, row = i % 2, i // 2
    x, y = 0.4 + col * 6.5, 1.25 + row * 2.9

    box(s3, x, y, 6.1, 2.65, RED_1)
    box(s3, x, y, 6.1, 0.07, RED_4)
    txt(s3, h, x + 0.18, y + 0.18, 5.7, 0.48, sz=13, bold=True, color=DARK_BG)
    txt(s3, b, x + 0.18, y + 0.74, 5.7, 1.65, sz=11, color=GRAY_7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 - HOW IT WORKS  (replaces "System Architecture")
# Non-technical: user journey in 3 steps + brief tech context at the bottom
# ══════════════════════════════════════════════════════════════════════════════
s4 = slide()
bg(s4, WHITE)
hdr(s4, "How It Works", "From raw recording to full meeting intelligence in a few clicks")

# Three step boxes
steps = [
    (INDIGO_5, "1. Upload",
     "Upload any meeting recording from your computer.\n\nSupports all common audio and video formats."),
    (VIOLET_5, "2. AI Processing",
     "The system listens, identifies each speaker, and reads through the entire conversation automatically."),
    (GREEN_5,  "3. Explore Insights",
     "A set of AI tools is instantly available:\nask questions, discover topics, and see how each person spoke."),
]

for i, (color, title, body) in enumerate(steps):
    x = 0.5 + i * 4.2

    box(s4, x, 1.25, 3.8, 4.45, INDIGO_1)
    box(s4, x, 1.25, 3.8, 0.08, color)

    # Step number circle approximation
    box(s4, x + 0.18, 1.42, 0.55, 0.55, color)
    txt(s4, str(i + 1), x + 0.18, 1.43, 0.55, 0.5,
        sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s4, title, x + 0.85, 1.5, 2.8, 0.48, sz=16, bold=True, color=DARK_BG)
    txt(s4, body,  x + 0.18, 2.12, 3.45, 2.5, sz=12, color=GRAY_7)

# Arrow connectors between steps
for i in range(2):
    ax = 4.3 + i * 4.2
    txt(s4, ">", ax, 2.9, 0.35, 0.6, sz=22, bold=True, color=INDIGO_5,
        align=PP_ALIGN.CENTER)

# Privacy & deployment note at the bottom
box(s4, 0.4, 5.88, 12.53, 1.4, INDIGO_1)
box(s4, 0.4, 5.88, 12.53, 0.06, INDIGO_5)

txt(s4, "Built with industry-standard tools. Runs entirely on your own machine - your data never leaves your computer.",
    0.65, 6.02, 12.1, 0.45, sz=12, bold=True, color=DARK_BG)

pillars = [
    (INDIGO_5, "Web Interface",
     "Clean, browser-based app.\nNo installation needed for the user."),
    (VIOLET_5, "Processing Engine",
     "Handles all requests between the\nweb app and the AI models."),
    (GREEN_5,  "Ready to Deploy",
     "One command to start the whole system.\nRuns on any machine with Docker."),
]
for i, (c, ptitle, pbody) in enumerate(pillars):
    px = 0.65 + i * 4.12
    box(s4, px, 6.55, 0.08, 0.6, c)
    txt(s4, ptitle, px + 0.18, 6.55, 3.7, 0.32, sz=10, bold=True, color=c)
    txt(s4, pbody,  px + 0.18, 6.88, 3.7, 0.38, sz=9,  color=GRAY_7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 - HOW THE AI WORKS, STEP BY STEP
# Fully custom-drawn flow: 4 pipeline steps (top) branching into 3 outputs (bottom)
# ══════════════════════════════════════════════════════════════════════════════
s5 = slide()
bg(s5, DARK_BG)
hdr(s5, "How the AI Works, Step by Step",
    "From raw recording to ready-to-use insights")

# ── Row 1: 4 pipeline step boxes ──────────────────────────────────────────────
# 4 boxes + 3 small arrow gaps across 12.73" usable width
BW   = 2.85   # box width
BH   = 2.05   # box height
BGAP = 0.41   # gap between boxes (arrow area)
BY   = 1.18   # top of row 1

step_boxes = [
    (INDIGO_5,
     "1.  Upload\nYour Recording",
     "Any audio or video file.\nMeetings, interviews, lectures\nand more."),
    (MED_BLUE,
     "2.  Speech\nto Text",
     "The AI listens to every word\nand writes it out in full\nwith timestamps."),
    (VIOLET_5,
     "3.  Speaker\nIdentification",
     "Each person's voice is\nrecognised and labelled\nautomatically."),
    (RGBColor(14, 165, 233),
     "4.  Content\nUnderstanding",
     "The AI reads the transcript\nto understand meaning,\ncontext, and key points."),
]

for i, (color, title, body) in enumerate(step_boxes):
    x = 0.3 + i * (BW + BGAP)

    # Box
    box(s5, x, BY,        BW, BH, DARK_CARD)
    box(s5, x, BY,        BW, 0.07, color)
    txt(s5, title, x + 0.15, BY + 0.15, BW - 0.25, 0.72,
        sz=11, bold=True, color=WHITE)
    txt(s5, body,  x + 0.15, BY + 0.96, BW - 0.25, 1.0,
        sz=10, color=INDIGO_3)

    # Arrow connector (horizontal bar + triangle head via thin boxes)
    if i < 3:
        ax = x + BW + 0.04
        ay = BY + BH / 2 - 0.02
        box(s5, ax, ay, BGAP - 0.08, 0.04, INDIGO_5)            # shaft
        box(s5, ax + BGAP - 0.14, ay - 0.1, 0.1, 0.24, INDIGO_5)  # head

# ── Connector: row 1 bottom → row 2 top ──────────────────────────────────────
SEP_Y = BY + BH + 0.08          # just below row 1
LABEL_Y = SEP_Y + 0.08
ROW2_Y  = LABEL_Y + 0.46

# Full-width thin bar
box(s5, 0.3, SEP_Y, 12.73, 0.05, INDIGO_5)

txt(s5, "Your insights, ready to explore:",
    0.3, LABEL_Y, 12.73, 0.38,
    sz=11, color=INDIGO_3, align=PP_ALIGN.CENTER)

# ── Row 2: 3 output feature boxes ─────────────────────────────────────────────
OBW  = 3.94   # output box width
OGAP = 0.50   # gap
OBH  = 7.38 - ROW2_Y - 0.12   # fill remaining height

output_boxes = [
    (INDIGO_5,
     "Ask the Meeting",
     "Type any question about the\nmeeting and the AI finds the\nanswer from the actual transcript.\n\nExample: \"What decisions were made?\"\nor \"What did Speaker A say?\""),
    (VIOLET_5,
     "Discover Topics",
     "The AI reads through all\ndiscussion segments and groups\nthem into clear, labelled themes.\n\nNo manual categorisation needed.\nSee what was really talked about."),
    (GREEN_5,
     "Read the Room",
     "See how each participant\ncame across during the meeting:\npositive, neutral, or negative.\n\nUseful for team health checks\nand post-meeting reflection."),
]

# Vertical drops from the separator bar down to each output box
out_centers = [0.3 + OBW / 2,
               0.3 + OBW + OGAP + OBW / 2,
               0.3 + 2 * (OBW + OGAP) + OBW / 2]

for cx in out_centers:
    box(s5, cx - 0.02, SEP_Y + 0.05, 0.04, ROW2_Y - SEP_Y - 0.05, INDIGO_5)

for i, (color, title, body) in enumerate(output_boxes):
    ox = 0.3 + i * (OBW + OGAP)

    box(s5, ox, ROW2_Y,        OBW, OBH, DARK_CARD)
    box(s5, ox, ROW2_Y,        OBW, 0.07, color)
    box(s5, ox, ROW2_Y + 0.07, OBW, 0.55, color)   # coloured title band
    txt(s5, title, ox + 0.15, ROW2_Y + 0.12, OBW - 0.25, 0.46,
        sz=13, bold=True, color=WHITE)
    txt(s5, body,  ox + 0.15, ROW2_Y + 0.72, OBW - 0.25, OBH - 0.82,
        sz=10.5, color=INDIGO_3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 - BUILT WITH
# Top: framework logos (Python, FastAPI, Streamlit, Docker, HuggingFace)
# Bottom: AI models end-to-end pipeline (Whisper → Pyannote → MiniLM → Qwen2 → RoBERTa)
# ══════════════════════════════════════════════════════════════════════════════
LOGO_DIR = r"D:\MeetRecall AI Improvement\documentations\logos"

s_tech = slide()
bg(s_tech, WHITE)
hdr(s_tech, "Built With",
    "Frameworks powering the platform and AI models used end-to-end")

# ─ Section label: Frameworks ─────────────────────────────────────────────────
txt(s_tech, "FRAMEWORKS & TOOLS",
    0.3, 1.13, 6, 0.3, sz=9, bold=True, color=INDIGO_5)

# ─ 5 logo cards ──────────────────────────────────────────────────────────────
LOGO_CARDS = [
    ("python.png",      "Python"),
    ("fastapi.png",     "FastAPI"),
    ("streamlit.png",   "Streamlit"),
    ("docker.png",      "Docker"),
    ("huggingface.png", "HuggingFace"),
]

LCW  = 2.32   # card width
LCH  = 2.15   # card height
LCGP = 0.27   # gap between cards
LCX0 = 0.30   # left start

for i, (fname, name) in enumerate(LOGO_CARDS):
    lx = LCX0 + i * (LCW + LCGP)
    ly = 1.48

    box(s_tech, lx, ly, LCW, LCH, INDIGO_1)
    box(s_tech, lx, ly, LCW, 0.06, INDIGO_5)

    logo_path = os.path.join(LOGO_DIR, fname)
    if os.path.exists(logo_path):
        # Centre logo horizontally; auto-height to preserve ratio
        s_tech.shapes.add_picture(
            logo_path,
            Inches(lx + 0.22), Inches(ly + 0.2),
            Inches(LCW - 0.44), Inches(1.38)
        )

    txt(s_tech, name,
        lx, ly + LCH - 0.44, LCW, 0.38,
        sz=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# ─ Separator ─────────────────────────────────────────────────────────────────
SEP2 = 1.48 + LCH + 0.1
box(s_tech, 0.3, SEP2, 12.73, 0.04, INDIGO_5)

txt(s_tech, "AI MODELS — END TO END",
    0.3, SEP2 + 0.1, 6, 0.3, sz=9, bold=True, color=INDIGO_5)

# ─ 5 model pipeline boxes ────────────────────────────────────────────────────
MODEL_BOXES = [
    (INDIGO_5,
     "Whisper",
     "Speech to Text",
     "Listens to the audio\nand writes out every\nword with timestamps"),
    (MED_BLUE,
     "Pyannote",
     "Speaker Identification",
     "Recognises each person's\nvoice and assigns a\nlabel automatically"),
    (VIOLET_5,
     "MiniLM",
     "Text Understanding",
     "Converts sentences into\nnumbers the AI can search\nand compare quickly"),
    (RGBColor(14, 165, 233),
     "Qwen2",
     "AI Reasoning",
     "Reads the transcript\nto answer questions and\nlabel discussion topics"),
    (GREEN_5,
     "RoBERTa",
     "Emotion Detection",
     "Classifies each speaker's\ntone as positive, neutral,\nor negative"),
]

MW   = 2.35    # model box width
MH   = 3.22    # model box height
MGAP = 0.26    # gap
MX0  = 0.30
MY   = SEP2 + 0.48

for i, (color, short, role, desc) in enumerate(MODEL_BOXES):
    mx = MX0 + i * (MW + MGAP)

    box(s_tech, mx, MY,        MW, MH, DARK_CARD)
    box(s_tech, mx, MY,        MW, 0.06, color)
    box(s_tech, mx, MY + 0.06, MW, 0.58, color)

    txt(s_tech, short, mx + 0.12, MY + 0.10, MW - 0.2, 0.44,
        sz=14, bold=True, color=WHITE)
    txt(s_tech, role,  mx + 0.12, MY + 0.72, MW - 0.2, 0.42,
        sz=9.5, bold=True, color=WHITE)
    txt(s_tech, desc,  mx + 0.12, MY + 1.22, MW - 0.2, 1.88,
        sz=10, color=INDIGO_3)

    # Arrow to next box (shaft + head)
    if i < 4:
        ax = mx + MW + 0.04
        ay = MY + MH / 2 - 0.02
        box(s_tech, ax, ay, MGAP - 0.1, 0.04, INDIGO_5)          # shaft
        box(s_tech, ax + MGAP - 0.16, ay - 0.1, 0.1, 0.24, INDIGO_5)  # head


# ══════════════════════════════════════════════════════════════════════════════
# DEMO SLIDES
# ══════════════════════════════════════════════════════════════════════════════
def demo(title, tag_color, caption_lines, img1, img2=None, img2_h=None):
    s = slide()
    bg(s, RGBColor(10, 8, 38))

    box(s, 0, 0, 13.33, 0.58, DARK_BG)
    box(s, 0, 0.58, 13.33, 0.04, tag_color)
    txt(s, title, 0.3, 0.08, 9, 0.46, sz=20, bold=True, color=WHITE)

    box(s, 0, 0.62, 2.35, 6.88, DARK_BG)
    y = 0.82
    for label, body in caption_lines:
        box(s, 0.12, y, 0.06, 0.38, tag_color)
        if label:
            txt(s, label, 0.25, y,        1.98, 0.36, sz=9.5, bold=True, color=WHITE)
            y += 0.42
            txt(s, body,  0.12, y,        2.15, 0.56, sz=9,   color=INDIGO_3)
            y += 0.68
        else:
            txt(s, body,  0.25, y,        1.98, 0.56, sz=9,   color=INDIGO_3)
            y += 0.72

    sx, sw, sy = 2.38, 10.88, 0.62
    if img2:
        img(s, img1, sx, sy,           sw, 3.46)
        box(s, sx, sy + 3.46, sw, 0.04, tag_color)
        if img2_h:
            img(s, img2, sx, sy + 3.5, sw, img2_h)
        else:
            img(s, img2, sx, sy + 3.5, sw)
    else:
        img(s, img1, sx, sy, sw)

    return s


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 - DEMO: UPLOAD & TRANSCRIPTION
# ══════════════════════════════════════════════════════════════════════════════
demo(
    title="Demo - Upload and Transcription",
    tag_color=INDIGO_5,
    caption_lines=[
        ("Upload",
         "Drag and drop any meeting file. The system accepts over 10 different audio and video formats."),
        ("Provider",
         "Choose between a fully private local option or a fast cloud option, depending on your preference."),
        ("Result",
         "A full transcript is produced with each speaker clearly identified and timestamped."),
        ("In this test",
         "24-minute recording, 7 speakers, 46 segments detected automatically"),
    ],
    img1="upload.png",
    img2="transcription_result.png",
    img2_h=3.32,
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 - DEMO: ASK YOUR MEETING  (RAG Chatbot)
# ══════════════════════════════════════════════════════════════════════════════
demo(
    title="Demo - Ask Your Meeting",
    tag_color=VIOLET_5,
    caption_lines=[
        ("How it works",
         "Type any question about the meeting and the AI finds the answer from the actual transcript."),
        ("Grounded answers",
         "Every response points back to the exact part of the meeting it came from."),
        ("What you can ask",
         "\"What decisions were made?\" \"What did Speaker A say about budgets?\" \"Summarise the meeting.\""),
        ("Language",
         "Works in English and Bahasa Indonesia"),
    ],
    img1="RAG_Chatbot.png",
    img2="more_interaction_chat.png",
    img2_h=3.32,
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 - DEMO: TOPIC DISCOVERY
# ══════════════════════════════════════════════════════════════════════════════
demo(
    title="Demo - Topic Discovery",
    tag_color=INDIGO_5,
    caption_lines=[
        ("What it does",
         "The AI reads through all discussion segments and automatically groups them into meaningful themes."),
        ("No manual work",
         "Topics and their labels are generated by the AI. No human categorisation required."),
        ("In this test",
         "5 distinct topics found across a 24-minute meeting with 46 discussion segments."),
        ("Use case",
         "Quickly understand what a long meeting was really about without re-watching it."),
    ],
    img1="topic_clustering.png",
    img2="detail_each_topic_clustering.png",
    img2_h=3.32,
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 - DEMO: READ THE ROOM  (Sentiment Analysis)
# ══════════════════════════════════════════════════════════════════════════════
demo(
    title="Demo - Read the Room",
    tag_color=GREEN_5,
    caption_lines=[
        ("What it does",
         "The AI analyses how each speaker came across throughout the meeting."),
        ("Per-person view",
         "Each participant gets their own positive, neutral, and negative breakdown."),
        ("In this test",
         "7 speakers analysed. Meeting tone was mostly neutral (59.8%), with 26.6% positive."),
        ("Use case",
         "Useful for team health checks, client call reviews, or post-meeting reflection."),
    ],
    img1="sentiment_analysis.png",
    img2="detail_per_speaker.png",
    img2_h=3.32,
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 - LET'S COLLABORATE + THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
s10 = slide()
bg(s10, DARK_BG)

box(s10, 8.62, 0, 4.71, 7.5, MED_BLUE)
box(s10, 0.45, 1.1, 0.08, 3.8, INDIGO_3)

txt(s10, "LET'S COLLABORATE!",
    0.7, 1.2, 7.6, 0.88, sz=35, bold=True, color=WHITE)

txt(s10,
    "Open to discussions, collaborations, and opportunities\n"
    "in Artificial Intelligence and Machine Learning.",
    0.7, 2.28, 7.2, 0.95, sz=14, color=INDIGO_3)

contacts = [
    ("Email",    "srilutfiyadwiy@gmail.com"),
    ("LinkedIn", "linkedin.com/in/sri-lutfiya-dwiyeni"),
    ("GitHub",   "github.com/rilufiyy"),
    ("Phone",    "+62 812-1457-7894"),
]
for i, (label, value) in enumerate(contacts):
    y = 3.5 + i * 0.72
    box(s10, 0.7, y + 0.04, 1.15, 0.44, MED_BLUE)
    txt(s10, label, 0.74, y + 0.06, 1.05, 0.36,
        sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s10, value, 1.98, y + 0.07, 5.9,  0.38, sz=12, color=WHITE)

box(s10, 0, 7.35, 13.33, 0.08, MED_BLUE)

txt(s10, "THANK YOU",
    8.85, 2.7, 4.2, 0.8, sz=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s10, "MeetRecall AI Improvement\nMeeting Intelligence Platform",
    8.85, 3.65, 4.2, 0.85, sz=12, color=WHITE, align=PP_ALIGN.CENTER)
s10.shapes.add_picture(
    os.path.join(DOCS, "main_interface.png"),
    Inches(8.88), Inches(4.72), Inches(4.05), Inches(2.0))

# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved : {OUT}")
print(f"Slides : {len(prs.slides)}")
